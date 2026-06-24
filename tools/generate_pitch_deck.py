from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "outputs"
PPTX_PATH = OUTPUT_DIR / "CyberShield_Police_Hackathon_Pitch.pptx"

NAVY = RGBColor(6, 19, 42)
NAVY_2 = RGBColor(10, 32, 63)
PANEL = RGBColor(13, 42, 79)
CYAN = RGBColor(34, 211, 238)
BLUE = RGBColor(59, 130, 246)
TEAL = RGBColor(45, 212, 191)
AMBER = RGBColor(251, 191, 36)
ROSE = RGBColor(251, 113, 133)
WHITE = RGBColor(248, 250, 252)
MUTED = RGBColor(188, 209, 231)
GRID = RGBColor(25, 72, 116)


def add_bg(slide, title: str, number: int, eyebrow: str = "CYBERSHIELD POLICE") -> None:
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    for x in [i * 0.7 for i in range(20)]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), 0, Inches(x), Inches(7.5))
        line.line.color.rgb = GRID
        line.line.transparency = 78
        line.line.width = Pt(0.5)
    for y in [i * 0.55 for i in range(15)]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0, Inches(y), Inches(13.333), Inches(y))
        line.line.color.rgb = GRID
        line.line.transparency = 82
        line.line.width = Pt(0.5)
    glow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.5), Inches(-1.3), Inches(4.8), Inches(4.8))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(13, 148, 136)
    glow.fill.transparency = 76
    glow.line.fill.background()

    add_text(slide, eyebrow, 0.55, 0.32, 4.0, 0.3, 9, CYAN, bold=True, caps=True)
    add_text(slide, title, 0.55, 0.68, 9.2, 0.62, 26, WHITE, bold=True)
    add_text(slide, f"{number:02d}", 12.35, 0.38, 0.42, 0.22, 8, MUTED, align=PP_ALIGN.RIGHT)


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int, color=WHITE, bold=False, align=PP_ALIGN.LEFT, caps=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text.upper() if caps else text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Aptos"
    return box


def card(slide, x, y, w, h, title, body, accent=CYAN, icon=""):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.2)
    shape.shadow.inherit = False
    if icon:
        add_text(slide, icon, x + 0.15, y + 0.12, 0.45, 0.36, 18, accent, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.22 + (0.38 if icon else 0), y + 0.13, w - 0.45, 0.32, 14, WHITE, bold=True)
    add_text(slide, body, x + 0.22, y + 0.62, w - 0.42, h - 0.72, 11, MUTED)
    return shape


def bullet_list(slide, bullets, x, y, w, h, size=17, color=MUTED):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Aptos"
        p.space_after = Pt(7)
    return box


def pill(slide, text, x, y, w, accent=CYAN):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(7, 47, 74)
    s.line.color.rgb = accent
    s.line.width = Pt(1)
    add_text(slide, text, x + 0.08, y + 0.08, w - 0.16, 0.18, 8, WHITE, bold=True, align=PP_ALIGN.CENTER)
    return s


def add_arrow(slide, x1, y1, x2, y2, color=CYAN):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(2.2)
    conn.line.end_arrowhead = True
    return conn


def add_bar_chart(slide, x, y, w, h):
    data = CategoryChartData()
    data.categories = ["OTP", "UPI", "KYC", "Phishing", "Job", "Loan"]
    data.add_series("Risk Signals", (84, 78, 75, 88, 62, 66))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.color.rgb = WHITE
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.color.rgb = MUTED
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.plots[0].data_labels.font.size = Pt(9)
    chart.plots[0].data_labels.font.color.rgb = WHITE
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = CYAN


def add_donut_chart(slide, x, y, w, h):
    data = CategoryChartData()
    data.categories = ["Critical", "High", "Medium", "Low"]
    data.add_series("Triage", (18, 28, 34, 20))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.font.size = Pt(10)
    chart.legend.font.color.rgb = WHITE
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.show_percentage = True
    chart.plots[0].data_labels.font.size = Pt(9)
    fills = [ROSE, AMBER, BLUE, TEAL]
    for idx, point in enumerate(chart.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = fills[idx]


def build_deck() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    slide = prs.slides.add_slide(blank)
    add_bg(slide, "CyberShield Police", 1, "PRAKASAM POLICE HACKATHON 2026 - MISSION Y4")
    add_text(slide, "AI-powered cybercrime triage for safer digital communities.", 0.65, 1.65, 8.2, 0.5, 23, CYAN, bold=True)
    add_text(slide, "Challenge 05 - Cybercrime Detection & Digital Fraud", 0.65, 2.18, 6.7, 0.36, 17, MUTED)
    card(slide, 0.65, 3.0, 3.1, 1.15, "Citizen Safety", "Report, scan, and receive clear guidance before acting.", TEAL, "01")
    card(slide, 4.0, 3.0, 3.1, 1.15, "Police Triage", "Risk score, priority, evidence metadata, and summaries.", CYAN, "02")
    card(slide, 7.35, 3.0, 3.1, 1.15, "Threat Intelligence", "Hotspots, synthetic indicators, and linked case patterns.", BLUE, "03")
    shield = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.HEXAGON, Inches(9.65), Inches(1.35), Inches(2.25), Inches(2.25))
    shield.fill.solid(); shield.fill.fore_color.rgb = RGBColor(7, 89, 133); shield.line.color.rgb = CYAN; shield.line.width = Pt(2.4)
    add_text(slide, "CS", 10.1, 1.94, 1.3, 0.58, 30, WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Team Lead / Developer: Anubrata Sarker", 0.65, 6.68, 5.4, 0.25, 11, MUTED)

    # 2
    slide = prs.slides.add_slide(blank); add_bg(slide, "Problem: Rising Cyber Fraud", 2)
    bullet_list(slide, ["Scams begin with small signals: SMS, WhatsApp, calls, links, QR codes", "Victims are pressured before they can verify", "Police receive scattered screenshots and incomplete reports", "Time-sensitive fraud needs fast prevention and triage"], 0.75, 1.55, 5.8, 3.3, 20)
    add_bar_chart(slide, 7.0, 1.55, 5.2, 3.3)
    card(slide, 0.75, 5.45, 11.5, 0.9, "Key Insight", "A useful prototype must help before loss, during reporting, and after police triage.", CYAN, "!")

    # 3
    slide = prs.slides.add_slide(blank); add_bg(slide, "Why Existing Systems Are Not Enough", 3)
    card(slide, 0.8, 1.45, 3.55, 1.35, "Reactive Reporting", "Many workflows begin after money or credentials are already lost.", ROSE, "A")
    card(slide, 4.85, 1.45, 3.55, 1.35, "Link-Only Detection Gap", "Plain SMS, calls, Telegram scripts, and OTP requests may have no URL.", AMBER, "B")
    card(slide, 8.9, 1.45, 3.55, 1.35, "Manual Triage Load", "Police need structured priority, category, and evidence context.", CYAN, "C")
    add_arrow(slide, 2.5, 3.55, 5.2, 3.55); add_arrow(slide, 6.7, 3.55, 9.35, 3.55)
    card(slide, 1.0, 4.35, 10.9, 1.35, "CyberShield Response", "Prevention-first citizen guidance plus explainable police triage in a single local prototype.", TEAL, "✓")

    # 4
    slide = prs.slides.add_slide(blank); add_bg(slide, "Our Solution: CyberShield Police", 4)
    bullet_list(slide, ["Unified citizen + police cyber safety platform", "Explainable risk score from 0 to 100", "Scam category, confidence, matched rules, and triage priority", "Gemini / multi-provider AI explanations with Local Safety Mode fallback", "Synthetic-data-only ethical prototype"], 0.8, 1.45, 5.9, 4.2, 18)
    for i, label in enumerate(["Message", "URL/UPI", "Transaction", "Case", "Hotspot"]):
        x = 7.0 + (i % 2) * 2.55
        y = 1.5 + (i // 2) * 1.15
        card(slide, x, y, 2.2, 0.8, label, "Risk signals", [CYAN, TEAL, AMBER, BLUE, ROSE][i], str(i + 1))
    add_text(slide, "Local detector is authoritative. AI only refines explanation text.", 7.0, 5.2, 4.8, 0.72, 18, CYAN, bold=True)

    # 5
    slide = prs.slides.add_slide(blank); add_bg(slide, "Five Core Modules", 5)
    modules = [
        ("Citizen Reporting", "Tracking ID, evidence metadata, safety guidance", TEAL),
        ("URL/UPI Validator", "Phishing, lookalike, UPI and message checks", CYAN),
        ("Transaction Anomaly", "Refund scams, fees, QR traps, repeated payments", AMBER),
        ("Threat Intelligence", "Hotspots, indicators, clusters, trend stream", BLUE),
        ("Police Dashboard", "Priority queue, case details, summaries, export", ROSE),
    ]
    for i, (t, b, c) in enumerate(modules):
        card(slide, 0.75 + i * 2.48, 2.0, 2.18, 2.15, t, b, c, f"{i+1}")
    add_text(slide, "One connected workflow, not separate demo screens.", 1.0, 5.25, 11.1, 0.45, 24, WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 6
    slide = prs.slides.add_slide(blank); add_bg(slide, "System Architecture", 6)
    boxes = [
        ("React + Vite UI", 0.8, 1.65, TEAL),
        ("FastAPI Routes", 3.35, 1.65, CYAN),
        ("Detection Engines", 5.9, 1.65, AMBER),
        ("SQLite + Synthetic Data", 8.45, 1.65, BLUE),
        ("Dashboards + Export", 10.35, 4.45, ROSE),
    ]
    for title, x, y, col in boxes:
        card(slide, x, y, 2.1, 1.0, title, "", col)
    add_arrow(slide, 2.9, 2.15, 3.35, 2.15); add_arrow(slide, 5.45, 2.15, 5.9, 2.15); add_arrow(slide, 8.0, 2.15, 8.45, 2.15)
    add_arrow(slide, 9.5, 2.75, 10.7, 4.45)
    card(slide, 1.05, 4.25, 7.7, 1.05, "Optional AI Provider Layer", "Gemini, OpenRouter, Groq, Hugging Face, or Local Safety Mode. API keys stay backend-only.", CYAN, "AI")
    add_arrow(slide, 5.9, 2.75, 5.6, 4.25)

    # 7
    slide = prs.slides.add_slide(blank); add_bg(slide, "AI Scam + Phishing + Screenshot Analysis", 7)
    bullet_list(slide, ["Text-only, URL-only, metadata-only, and hybrid detection", "OTP/PIN/CVV/password, KYC threats, UPI refunds, fake jobs, loans, QR scams", "Trusted-domain safeguards for false-positive prevention", "Screenshot analyzer uses metadata + pasted visible text, not private image storage"], 0.75, 1.45, 6.3, 3.5, 17)
    add_donut_chart(slide, 7.45, 1.5, 4.6, 3.4)
    card(slide, 1.0, 5.5, 11.1, 0.8, "Explainable Output", "Risk score • Risk level • Category • Confidence • Matched patterns • Police triage priority", CYAN, "✓")

    # 8
    slide = prs.slides.add_slide(blank); add_bg(slide, "CyberDost + Multilingual Citizen Support", 8)
    card(slide, 0.8, 1.45, 4.1, 1.35, "CyberDost Chatbot", "Practical, calm safety guidance with quick actions for UPI, phishing, digital arrest, reporting, and 1930.", CYAN, "BOT")
    card(slide, 5.1, 1.45, 3.4, 1.35, "Language Access", "English, Telugu, Hindi, Bengali, Tamil, Kannada, Malayalam, Marathi.", TEAL, "LANG")
    card(slide, 8.7, 1.45, 3.65, 1.35, "Local Safety Mode", "Guidance continues even when AI provider keys are absent or busy.", AMBER, "SAFE")
    langs = ["EN", "TE", "HI", "BN", "TA", "KN", "ML", "MR"]
    for i, lang in enumerate(langs):
        pill(slide, lang, 1.0 + i * 1.45, 4.15, 0.85, CYAN if i < 2 else BLUE)
    add_text(slide, "Citizen-first design: prevention guidance before money, OTP, or credentials are shared.", 1.0, 5.25, 11.0, 0.7, 22, WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 9
    slide = prs.slides.add_slide(blank); add_bg(slide, "Transaction Monitoring + Hotspot Intelligence", 9)
    add_bar_chart(slide, 0.9, 1.45, 5.2, 3.3)
    bullet_list(slide, ["Fake refund and chargeback patterns", "Loan processing fee and prize fee traps", "Job registration fees and first-time receivers", "Late-night transfers, repeated small payments, senior/student context", "Prakasam/AP synthetic hotspot and threat stream analytics"], 6.75, 1.45, 5.6, 3.6, 17)
    card(slide, 0.9, 5.35, 11.4, 0.85, "Admin/SP View", "High-risk transaction count, synthetic monitored amount, hotspot trend, and priority locations.", TEAL, "SP")

    # 10
    slide = prs.slides.add_slide(blank); add_bg(slide, "Police Workflow: Report → Analyze → Triage → Investigate → Warn Public", 10)
    steps = [("Report", "Citizen submits message, URL, transaction, evidence metadata"), ("Analyze", "Detector scores risk and extracts factors"), ("Triage", "Priority and police summary generated"), ("Investigate", "Officer updates status and notes"), ("Warn", "Hotspot/threat insights support awareness")]
    for i, (t, b) in enumerate(steps):
        x = 0.6 + i * 2.5
        card(slide, x, 2.0, 2.05, 1.65, t, b, [TEAL, CYAN, AMBER, BLUE, ROSE][i], str(i + 1))
        if i < 4:
            add_arrow(slide, x + 2.06, 2.82, x + 2.42, 2.82)
    card(slide, 1.0, 5.1, 11.1, 0.9, "Operational Value", "Structured cases, evidence metadata, FIR-style summary aid, PDF export, and dashboard analytics.", CYAN, "OPS")

    # 11
    slide = prs.slides.add_slide(blank); add_bg(slide, "Practical Impact, Scalability, Privacy & Ethics", 11)
    card(slide, 0.8, 1.5, 3.5, 1.35, "Impact", "Prevents fraud earlier and helps citizens choose safe next steps.", TEAL, "I")
    card(slide, 4.9, 1.5, 3.5, 1.35, "Scalability", "Can move to PostgreSQL, queues, OCR, audit logs, and official integrations.", CYAN, "S")
    card(slide, 9.0, 1.5, 3.5, 1.35, "Ethics", "Synthetic data only. No real victim data, police data, UPI IDs, or private screenshots.", AMBER, "E")
    bullet_list(slide, ["Local-first prototype runs without paid APIs", "AI does not override risk score or category", "Backend-only keys; .env ignored by Git", "Translations are demo-level and should be reviewed before real deployment"], 1.0, 4.0, 10.9, 1.7, 18)

    # 12
    slide = prs.slides.add_slide(blank); add_bg(slide, "Future Scope + Team", 12)
    roadmap = [("Official Integrations", "1930 / cybercrime portal workflows"), ("Secure Evidence", "Consent-based OCR and evidence vault"), ("District Intelligence", "Cross-district clustering and heatmaps"), ("Production Governance", "Audit logs, RBAC, native-language review")]
    for i, (t, b) in enumerate(roadmap):
        card(slide, 0.85 + (i % 2) * 5.8, 1.55 + (i // 2) * 1.6, 5.2, 1.05, t, b, [CYAN, TEAL, BLUE, AMBER][i], str(i + 1))
    card(slide, 1.2, 5.35, 10.7, 0.9, "Team", "Team Lead / Developer: Anubrata Sarker • Member 2: Placeholder Member • Member 3: Placeholder Member", CYAN, "TEAM")
    add_text(slide, "CyberShield Police: prevention-first, explainable, ethical, and demo-ready.", 1.0, 6.55, 11.2, 0.35, 17, WHITE, bold=True, align=PP_ALIGN.CENTER)

    OUTPUT_DIR.mkdir(exist_ok=True)
    prs.save(PPTX_PATH)


if __name__ == "__main__":
    build_deck()
    print(PPTX_PATH)
